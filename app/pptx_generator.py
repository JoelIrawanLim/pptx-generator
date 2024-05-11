#!/usr/bin/env python3

from pptx import Presentation
from datetime import datetime, timedelta
import os
import sys
import yaml
from importlib import resources
from rich.console import Console
from rich.table import Table
# from textual.app import App, ComposeResult
# from textual.widgets import Button, Digits

# Puts values on all the gloabal variables
def initialize():
    global console
    global data 
    global prs 
    global output_presentation
    global song_1_id
    global song_2_id
    global song_3_id
    console = Console()
    # Gets the path of the script. As the script might be downloaded in different locations. 
    script_path = os.path.abspath(__file__)
    script_dir = os.path.split(script_path)[0]
    # Comines the path where the script is locaed with the relative position of the yaml file
    with open(os.path.join(script_dir, "data/data.yml"), 'r') as file:
        data = yaml.safe_load(file)

    # Combines the path of where the script is located with the relative location of where the template Presentation is. 
    pptx_path = os.path.join(script_dir, "data/Remaja_template.pptx")
    prs = Presentation(pptx_path)
    # The output presentation's path and name. This is a placeholder so that it still generates a presentation when no name is provided. 
    output_presentation = "presentation.pptx"
    # Placeholder Song_ID
    song_1_id = '003'
    song_2_id = '001'
    song_3_id = '004'

# Gets the date of the next wednesday. 
def get_next_weekday(startdate, weekday):
    """
    @startdate: given date, in format '2013-05-25'
    @weekday: week day as a integer, between 0 (Monday) to 6 (Sunday)
    """
    d = datetime.strptime(startdate, '%d %B %Y')
    t = timedelta((7 + weekday - d.weekday()) % 7)
    return (d + t).strftime('%d %B %Y')

# Gets the object of the song based on the song_id. Song id k1 will give the object of k1
def id_to_song():
    global song_1
    global song_2 
    global song_3
    song_1 = data[song_1_id]
    song_2 = data[song_2_id]
    song_3 = data[song_3_id]

# Converts verse_number into a string, which is useable for the script
def generate_verse_number(verse_num): 
    if isinstance(verse_num, int):
        return str(verse_num)
    elif isinstance(verse_num, str):
        return verse_num
    else: 
        return 0

# Adds a black blank slide. 
def add_black_slide():
    prs.slides.add_slide(prs.slide_layouts[7])

# Adds the first silde which says "Welcome to Remaja"
def add_welcome_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    date_text = slide.placeholders[10]
    church_name = slide.placeholders[11]
    title.text = "Welcome to Remaja"
    date_text.text = get_next_weekday(datetime.today().strftime('%d %B %Y'), 5)
    church_name.text = "Reformed Evangelical Church Singapore"
    add_black_slide()

# Adds the title slide for the first song, with the proper template and color
def add_first_song_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = song_1['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = author_find(song_1_id)
    number_text.text = "1"

# Adds the title slide for the second song, with the proper template and color
def add_second_song_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = song_2['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = author_find(song_2_id)
    number_text.text = "2"

# Adds the title slide for the third song, with the proper template and color
def add_third_song_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title = slide.shapes.title
    title.text = song_3['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = author_find(song_3_id)
    number_text.text = "3"
    
# Adds the lyric slides for the first song, with the proper template and color
def add_first_song_lyrics_slide(verse_number):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    title = slide.shapes.title
    title.text = song_1['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = author_find(song_1_id)
    if (verse_number+1 > len(song_1['lyrics'])):
        console.print(f"[red bold] An error occured. If you are seeing this, send a bug report. The song is {song_1_id}")
        return 0
    lyrics_text.text = song_1['lyrics'][verse_number][0]
    number_text.text = generate_verse_number(song_1['verse_number'][verse_number])
    for i in range(1, len(song_1['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_1['lyrics'][verse_number][i]

# Adds the lyric slides for the second song, with the proper template and color
def add_second_song_lyrics_slide(verse_number):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = song_2['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = author_find(song_2_id)
    number_text.text = generate_verse_number(song_2['verse_number'][verse_number])
    if (verse_number+1 > len(song_2['lyrics'])):
        console.print(f"[red bold] An error occured. If you are seeing this, send a bug report. The song is {song_2_id}")
        return 0
    lyrics_text.text = song_2['lyrics'][verse_number][0]
    for i in range(1, len(song_2['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_2['lyrics'][verse_number][i]

# Adds the lyric slides for the third song, with the proper template and color
def add_third_song_lyrics_slide(verse_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.title
    title.text = song_3['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = author_find(song_3_id)
    number_text.text = generate_verse_number(song_3['verse_number'][verse_number])
    if (verse_number+1 > len(song_3['lyrics'])):
        console.print(f"[red bold] An error occured. If you are seeing this, send a bug report. The song is {song_3_id}")
        return 0
    lyrics_text.text = song_3['lyrics'][verse_number][0]
    for i in range(1, len(song_3['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_3['lyrics'][verse_number][i]    

# Add multiple slides in the announcement time. 
# Includes the slide which says announcements, 
#  A blank white slide for placing images of the announcement,
# The birthday song,
# And the "See you next week!" Slide
def add_announcements():
    prs.slides.add_slide(prs.slide_layouts[8])
    prs.slides.add_slide(prs.slide_layouts[9])
    birthday = prs.slides.add_slide(prs.slide_layouts[10])
    birthday.shapes.title.text = "Birthday Song"
    prs.slides.add_slide(prs.slide_layouts[11])

# Function which run the above functions for generating individual slides for the first slide. 
# Runs the title slide, the appropiate number of lyric slides, and a blank black slide
def generate_first_song():
    add_first_song_title_slide()
    for i in range(0, len(song_1['verse_number'])):
        add_first_song_lyrics_slide(i)
    add_black_slide()

# Function which run the above functions for generating individual slides for the second slide. 
# Runs the title slide, the appropiate number of lyric slides, and a blank black slide
def generate_second_song():
    add_second_song_title_slide()
    for i in range(0, len(song_2['verse_number'])):
        add_second_song_lyrics_slide(i)
    add_black_slide()

# Function which run the above functions for generating individual slides for the third slide. 
# Runs the title slide, the appropiate number of lyric slides, and a blank black slide
def generate_third_song():
    add_third_song_title_slide()
    for i in range(0, len(song_3['verse_number'])):
        add_third_song_lyrics_slide(i)
    add_black_slide()

# Formats the text to be lowercase, removes any _ and -
# Checks if the first 3 letters are "kri", and if so, removes any 0 directly behind the kri text (K0000010 becones K10)
# If it is only a k, it removes any 0 behind kri and removes it. And replaces k with kri. 
def text_formater(text):
    text = text.lower().replace("_", "").replace("-", "")
  #  if text[0] == "k":
  #      if text[1] == "r" and text[2] == "i":
   #         text_temporary = text.replace("kri", "")
    #        while text_temporary[0] == "0": 
 #               text_temporary = text_temporary[1:]
#            return "kri" + text_temporary
#        else:
  #          text_temporary = text.replace("k", "")
  #          while text_temporary[0] == "0": 
   #             text_temporary = text_temporary[1:]
   #         return "kri" + text_temporary
  #  else:
    return text

# Checks if the song id is a valid song id inside the song-list
def song_id():
    global song_1_id
    global song_2_id
    global song_3_id
    global Fail
    argv1=text_formater(sys.argv[1])
    argv2=text_formater(sys.argv[2])
    argv3=text_formater(sys.argv[3])
    Fail = 0 
    if data.get(argv1):
        song_1_id = argv1
        Fail += 1
    else: 
        song_1_id = argv1
        print(f"{song_1_id} is not a valid song id" )
    if data.get(argv2):
        song_2_id = argv2
        Fail += 1
    else: 
        song_2_id = argv2
        print(f"{song_2_id} is not a valid song id" )
    if data.get(argv3):
        song_3_id = argv3
        Fail += 1
    else: 
        song_3_id = argv3
        print(f"{song_3_id} is not a valid song id" )
    if Fail != 3: 
        sys.exit()

def author_find(key):
    text = data[key]['author'][0]
    for i in range(1, len(data[key]['author'])):
        text = f"{text}, {data[key]['author'][i]}"
    return text
        

#----------------------------- MAIN FUNCTION -----------------------------------------------------------------

# This is the main function, which is run when activating pptx-generator
def main():
    global output_presentation
    initialize()
    if len(sys.argv) == 4:
        # song_1_id = sys.argv[1]
        # song_2_id = sys.argv[2]
        # song_3_id = sys.argv[3]
        song_id()
        id_to_song()
        add_welcome_slide()
        generate_first_song()
        generate_second_song()
        generate_third_song() 
        add_announcements() 
        prs.save(output_presentation)
        console.print(f"[green]Done! Saved [/]{output_presentation}")
    elif len(sys.argv) > 4: 
        # song_1_id = sys.argv[1]
        # song_2_id = sys.argv[2]
        # song_3_id = sys.argv[3]
        song_id()
        if sys.argv[4].endswith(('.pptx')):
            output_presentation = sys.argv[4]
        else:
            output_presentation = sys.argv[4] + ".pptx"
        id_to_song()
        add_welcome_slide()
        generate_first_song()
        generate_second_song()
        generate_third_song() 
        add_announcements() 
        prs.save(output_presentation)
        console.print(f"[green]Done! Saved [/]{output_presentation}")
       
    else:
        console.print("[red bold]Not enough arguments provided. Requires at least 3 values")

# This is the search function, which is run when doing pptx-generator-search
def search():
    global output_presentation
    initialize()
    table = Table(title="Search Results")
    table.add_column("PPTX-Address", style="cyan")
    table.add_column("Title", style="blue")
    table.add_column("Author", style="green")
    if len(sys.argv) == 3: 
        found = 0
        if sys.argv[2] == "-a":
            addresses = []
            for key in data:
                for index_of_all_authors in range(0, len(data[key]['author'])):
                    for index_of_all_words_of_authors in range(0, len(data[key]['author'][index_of_all_authors].split(" "))):
                        for index_of_all_userinput_words in range(0, len(sys.argv[1].split(" "))):
                            if sys.argv[1].split(" ")[index_of_all_userinput_words].lower().replace(" ", "").replace("?", "").replace("!", "").replace(".","") == data[key]['author'][index_of_all_authors].split(" ")[index_of_all_words_of_authors].lower().replace(" ", "").replace("?", "").replace("!", "").replace(".",""):
                                already_done = False
                                for done_addresses in range(0, len(addresses)):
                                    if addresses[done_addresses] == key: 
                                        already_done = True
                                if already_done == False:
                                    table.add_row(key, data[key]['title'], author_find(key))
                                     # print( f"PPTX-Adress: {key}, Title: {data[key]['title']}, Author: {data[key]['author']}")
                                    found = 1
                                    addresses.append(key)
                                    

            if found != 1: 
                console.print(f"There is no author named [red underline bold]{sys.argv[1]}[/] in our song list.")
            else:
                console.print(table)
        elif sys.argv[2] == "-t":
            for key in data:
                if sys.argv[1].lower().replace(" ", "").replace("?", "").replace("!", "").replace(".","").replace(",","").replace(";","").replace(":","").replace("'","") == data[key]['title'].lower().replace(" ", "").replace("?", "").replace("!", "").replace(".","").replace(",","").replace(";","").replace(":","").replace("'",""):
                    table.add_row(key, data[key]['title'], author_find(key))
                    # print( f"PPTX-Adress: {key}, Title: {data[key]['title']}, Author: {data[key]['author']}")
                    found = 1
            if found != 1: 
                console.print(f"There is no song named [red underline bold]{sys.argv[1]}[/] in our song list.")
            else:
                console.print(table)
        elif sys.argv[2] == "-k":
            if sys.argv[1].isdigit():
                for key in data:
                    if int(sys.argv[1]) == data[key]['kri_number']:
                        table.add_row(key, data[key]['title'], author_find(key))
                        # print( f"PPTX-Adress: {key}, Title: {data[key]['title']}, Author: {data[key]['author']}")
                        found = 1
                if found != 1: 
                    console.print(f"There is no [red underline bold] KRI {sys.argv[1]}[/] in our song list.")
                else:
                    console.print(table)
            else:
                console.print("[red bold]That is not a valid KRI number")
        else: 
            console.print("[red bold]That is not a valid flag. pptx-generator-search only accepts -k, -a, and -t flags")
    else: 
        print('pptx-generator-search requires 2 values. \nThe first is the author, title, or KRI number of the song surrounded in double quotes ("), \nNext is the flag. -a for author, -t for title, and -k for kri number\nExample: pptx-generator-search "Grace Alone" -t')

# Checks if the function is run using the run button in VScode, and then activates the main function. 
if __name__ == "__main__":
    # pptx_generator().run()
    main()

# Debug Function which generates a pptx of all the songs in all slide forms
def debug():
    global song_1
    global song_2 
    global song_3
    global output_presentation
    global song_1_id
    global song_2_id
    global song_3_id
    initialize()
    # song_1_id = sys.argv[1]
    # song_2_id = sys.argv[2]
    # song_3_id = sys.argv[3]
    add_welcome_slide()
    for number in data:
        if number != "empty":
            song_1_id = number
            song_2_id = number 
            song_3_id = number
            song_1 = data[number]
            song_2 = data[number]
            song_3 = data[number]
            generate_first_song()
            generate_second_song()
            generate_third_song() 
    add_announcements() 
    prs.save(output_presentation)
    console.print(f"[green]Done! Saved [/]{output_presentation}")
