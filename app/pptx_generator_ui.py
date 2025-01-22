from datetime import datetime, timedelta
from pptx import Presentation
import os
import yaml
import sys

# global variables for search
search_query = ""  #user input
search_option = "" # -a, -t or -k
searchresult_key = []
searchresult_title= []
searchresult_author= []
searchresult_error= ""
searchresult_invalid= ""

# global variables for generate
first_song = [] 
second_song = []
third_song = []
first_song_id = ""
second_song_id = ""
third_song_id = ""
presentation_title = ""
output_presentation = ""

# main app
# Gets the date of the next wednesday.
def get_next_weekday(startdate, weekday):
   """
    @startdate: given date, in format '2013-05-25'
    @weekday: week day as a integer, between 0 (Monday) to 6 (Sunday)
    """
   d = datetime.strptime(startdate, '%d %B %Y')
   t = timedelta((7 + weekday - d.weekday()) % 7)
   return (d + t).strftime('%d %B %Y')


def initialize():
   global data
   global prs
   global script_dir
   # Gets the path of the script. As the script might be downloaded in different locations.
   script_path = os.path.abspath(__file__)
   script_dir = os.path.split(script_path)[0]
   # Comines the path where the script is locaed with the relative position of the yaml file
   with open(os.path.join(script_dir, "data/data.yml"), 'r') as file:
      data = yaml.safe_load(file)
   # Combines the path of where the script is located with the relative location of where the template Presentation is.
   pptx_path = os.path.join(script_dir, "data/Remaja_template.pptx")
   prs = Presentation(pptx_path)


def id_to_song(first_song_id, second_song_id, third_song_id):
   global first_song
   global second_song
   global third_song
   first_song = data[first_song_id]
   second_song = data[second_song_id]
   third_song = data[third_song_id]


def generate_verse_number(verse_num):
   if isinstance(verse_num, int):
      return str(verse_num)
   elif isinstance(verse_num, str):
      return verse_num
   else:
      return 0


# Adds a black blank slide.
def add_black_slide():
   prs.slides.add_slide(prs.slide_layouts[7])


# Adds the first silde which says "Welcome to Remaja"
def add_welcome_slide():
   slide = prs.slides.add_slide(prs.slide_layouts[0])
   title = slide.shapes.title
   date_text = slide.placeholders[10]
   church_name = slide.placeholders[11]
   title.text = "Welcome to Remaja"
   date_text.text = get_next_weekday(datetime.today().strftime('%d %B %Y'), 5)
   church_name.text = "Reformed Evangelical Church Singapore"
   add_black_slide()


def find_author(key):
   text = key['author'][0]
   for i in range(1, len(key['author'])):
      text = f"{text}, {key['author'][i]}"
   return text


# Adds the title slide for the first song, with the proper template and color
def add_first_song_title_slide():
   global first_song
   slide = prs.slides.add_slide(prs.slide_layouts[1])
   title = slide.shapes.title
   title.text = first_song['title']
   number_text = slide.placeholders[10]
   author_text = slide.placeholders[11]
   author_text.text = find_author(first_song)
   number_text.text = "1"


# Adds the title slide for the second song, with the proper template and color
def add_second_song_title_slide():
   slide = prs.slides.add_slide(prs.slide_layouts[2])
   title = slide.shapes.title
   title.text = second_song['title']
   number_text = slide.placeholders[10]
   author_text = slide.placeholders[11]
   author_text.text = find_author(second_song)
   number_text.text = "2"


# Adds the title slide for the third song, with the proper template and color
def add_third_song_title_slide():
   slide = prs.slides.add_slide(prs.slide_layouts[3])
   title = slide.shapes.title
   title.text = third_song['title']
   number_text = slide.placeholders[10]
   author_text = slide.placeholders[11]
   author_text.text = find_author(third_song)
   number_text.text = "3"


# Adds the lyric slides for the first song, with the proper template and color
# Inputs is verse number, which is an array of the numbers written at the bottom of the slides.
def add_first_song_lyrics_slide(verse_number):
   slide = prs.slides.add_slide(prs.slide_layouts[4])
   title = slide.shapes.title
   title.text = first_song['title']
   author_text = slide.placeholders[11]
   number_text = slide.placeholders[13]
   lyrics_text = slide.placeholders[12].text_frame
   author_text.text = find_author(first_song)
   if (verse_number + 1 > len(first_song['lyrics'])):
      # Fix this afterward
      print("An Error Occured.")
      return 0
   lyrics_text.text = first_song['lyrics'][verse_number][0]
   number_text.text = generate_verse_number(
       first_song['verse_number'][verse_number])
   for i in range(1, len(first_song['lyrics'][verse_number])):
      p = lyrics_text.add_paragraph()
      p.text = first_song['lyrics'][verse_number][i]


# Adds the lyric slides for the second song, with the proper template and color
# Inputs is verse number, which is an array of the numbers written at the bottom of the slides.
def add_second_song_lyrics_slide(verse_number):
   slide = prs.slides.add_slide(prs.slide_layouts[5])
   title = slide.shapes.title
   title.text = second_song['title']
   author_text = slide.placeholders[11]
   number_text = slide.placeholders[13]
   lyrics_text = slide.placeholders[12].text_frame
   author_text.text = find_author(second_song)
   number_text.text = generate_verse_number(
       second_song['verse_number'][verse_number])
   if (verse_number + 1 > len(second_song['lyrics'])):
      # Fix this afterward
      print("An Error Occured.")
      return 0
   lyrics_text.text = second_song['lyrics'][verse_number][0]
   for i in range(1, len(second_song['lyrics'][verse_number])):
      p = lyrics_text.add_paragraph()
      p.text = second_song['lyrics'][verse_number][i]


# Adds the lyric slides for the third song, with the proper template and color
# Inputs is verse number, which is an array of the numbers written at the bottom of the slides.
def add_third_song_lyrics_slide(verse_number):
   slide = prs.slides.add_slide(prs.slide_layouts[6])
   title = slide.shapes.title
   title.text = third_song['title']
   author_text = slide.placeholders[11]
   number_text = slide.placeholders[13]
   lyrics_text = slide.placeholders[12].text_frame
   author_text.text = find_author(third_song)
   number_text.text = generate_verse_number(
       third_song['verse_number'][verse_number])
   if (verse_number + 1 > len(third_song['lyrics'])):
      # Fix this afterward
      print("An Error Occured.")
      return 0
   lyrics_text.text = third_song['lyrics'][verse_number][0]
   for i in range(1, len(third_song['lyrics'][verse_number])):
      p = lyrics_text.add_paragraph()
      p.text = third_song['lyrics'][verse_number][i]


# Add multiple slides in the announcement time.
# Includes the slide which says announcements,
#  A blank white slide for placing images of the announcement,
# The birthday song,
# And the "See you next week!" Slide
# No inputs
def add_announcements():
   prs.slides.add_slide(prs.slide_layouts[8])
   prs.slides.add_slide(prs.slide_layouts[9])
   birthday = prs.slides.add_slide(prs.slide_layouts[10])
   birthday.shapes.title.text = "Birthday Song"
   birthday_lyrics = birthday.placeholders[10].text_frame
   birthday_lyrics_paragraph = birthday_lyrics.paragraphs[0]
   birthday_lyrics_paragraph.text = "Selamat ulang tahun,"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Kami ucapkan padamu."
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Selamat hari jadi,"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Tuhan Yesus memberkati."
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Kami s'lalu berdoa"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Agar kau tetap setia"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Pada Yesus Kristus, Tuhan dan Raja kita, oh."
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Kami mengucap syukur"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Tuhan t'lah pimpin langkahmu"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   birthday_lyrics_paragraph.text = "Padamu [insert name], selamat ulang tahun!"
   birthday_lyrics_paragraph = birthday_lyrics.add_paragraph()
   prs.slides.add_slide(prs.slide_layouts[11])


def generate_first_song():
   add_first_song_title_slide()
   for i in range(0, len(first_song['verse_number'])):
      add_first_song_lyrics_slide(i)
   add_black_slide()


def generate_second_song():
   add_second_song_title_slide()
   for i in range(0, len(second_song['verse_number'])):
      add_second_song_lyrics_slide(i)
   add_black_slide()


def generate_third_song():
   add_third_song_title_slide()
   for i in range(0, len(third_song['verse_number'])):
      add_third_song_lyrics_slide(i)
   add_black_slide()


def text_formatter(text):
   text = text.lower().replace("_", "").replace("-", "")
   return text


# Checks if the song id is a valid song id inside the song-list
def song_id_checker(first_song_id, second_song_id, third_song_id):
   Fail = 0
   if data.get(first_song_id):
      Fail += 1
   else:
      print(f"{first_song_id} is not a valid song id")
   if data.get(second_song_id):
      Fail += 1
   else:
      print(f"{second_song_id} is not a valid song id")
   if data.get(third_song_id):
      Fail += 1
   else:
      print(f"{third_song_id} is not a valid song id")
   if Fail != 3:
      sys.exit()


def main(first_song_id, second_song_id, third_song_id, presentation_title):
   global output_presentation
   initialize()
   song_id_checker(text_formatter(first_song_id), text_formatter(second_song_id), text_formatter(third_song_id))
   id_to_song(first_song_id, second_song_id, third_song_id)
   add_welcome_slide()
   generate_first_song()
   generate_second_song()
   generate_third_song()
   add_announcements()
   if presentation_title.endswith(('.pptx')):
      output_presentation = presentation_title
   else:
      output_presentation = presentation_title + ".pptx"
   prs.save(script_dir + "/"  + output_presentation)


def search(search_query, search_option):
   global searchresult_key
   global searchresult_title 
   global searchresult_author
   global searchresult_invalid 
   global searchresult_error
   searchresult_key = []
   searchresult_title = []
   searchresult_author = []
   searchresult_invalid = []
   searchresult_error = []
   initialize()
   found = 0
   if search_option == "-a":
      addresses = []
      for key in data:
         for index_of_all_authors in range(0, len(data[key]['author'])):
            for index_of_all_words_of_authors in range(
                0, len(data[key]['author'][index_of_all_authors].split(" "))):
               for index_of_all_userinput_words in range(0, len(search_query.split(" "))):
                  if search_query.split( " ")[index_of_all_userinput_words].lower().replace(" ", "").replace("?", "").replace("!", "").replace(
                      ".", "") == data[key]['author'][index_of_all_authors].split(" ")[index_of_all_words_of_authors].lower().replace(" ","").replace("?", "").replace("!","").replace(".", ""):
                     already_done = False
                     for done_addresses in range(0, len(addresses)):
                        if addresses[done_addresses] == key:
                           already_done = True
                     if not already_done:
                        searchresult_key.append(key)
                        searchresult_title.append(data[key]['title'])
                        searchresult_author.append(data[key]['author'])
                        found = 1
                        addresses.append(key)
      if found != 1:
         searchresult_error = f"There is no author named {search_query} in our song list."

   elif search_option == "-t":
      found = 0
      for key in data:
         if search_query.lower().replace(" ", "").replace("?", "").replace(
             "!",
             "").replace(".", "").replace(",", "").replace(";", "").replace(
                 ":",
                 "").replace("'", "") == data[key]['title'].lower().replace(
                     " ", "").replace("?", "").replace("!", "").replace(
                         ".", "").replace(",", "").replace(";", "").replace(
                             ":", "").replace("'", ""):
            searchresult_key.append(key)
            searchresult_title.append(data[key]['title'])
            searchresult_author.append(data[key]['author'])
            searchresult_error = ""
            searchresult_invalid =""
            found = 1
      if found != 1:
         searchresult_error = f"There is no song named {search_query} in our song list."
         
   elif search_option == "-k":
      if search_query.isdigit():
         for key in data:
            if int(search_query) == data[key]['kri_number']:
               searchresult_key.append(key)
               searchresult_title.append(data[key]['title'])
               searchresult_author.append(data[key]['author'])
               searchresult_error = ""
               searchresult_invalid =""
               found = 1
         if found != 1:
            searchresult_error = f"There is no kri{search_query} in our song list."
      else:
         searchresult_invalid = ("Invalid input, only accepts numbers.")

   


if __name__ == "__main__":
   pass
