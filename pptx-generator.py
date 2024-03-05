from pptx import Presentation
from datetime import datetime, timedelta
import os
import yaml
with open('data.yml', 'r') as file:
   data = yaml.safe_load(file)

def get_next_weekday(startdate, weekday):
    """
    @startdate: given date, in format '2013-05-25'
    @weekday: week day as a integer, between 0 (Monday) to 6 (Sunday)
    """
    d = datetime.strptime(startdate, '%d %B %Y')
    t = timedelta((7 + weekday - d.weekday()) % 7)
    return (d + t).strftime('%d %B %Y')

class song:
   def __init__(self, title, author, verse_number, lyrics):
      self.title = title
      self.author = author
      self.verse_number = verse_number
      self.lyrics = lyrics

script_path = os.path.abspath(__file__)
script_dir = os.path.split(script_path)[0]
prs = Presentation(os.path.join(script_dir, "presentations/Remaja_template.pptx"))
output_presentation = './presentations/test4.pptx'
# ------------------- Song ID -------------------------------
song_1_id = '003'
song_2_id = '001'
song_3_id = '004'

def id_to_song():
    global song_1
    global song_2 
    global song_3
    song_1 = data[song_1_id]
    song_2 = data[song_2_id]
    song_3 = data[song_3_id]

def generate_verse_number(verse_num): 
   if isinstance(verse_num, int):
      return str(verse_num)
   elif isinstance(verse_num, str):
      return verse_num
   else: 
      return 0
def add_welcome_slide():
    print("Welcome Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    date_text = slide.placeholders[10]
    church_name = slide.placeholders[11]
    title.text = "Welcome to Remaja"
    # date_text.text = datetime.strptime(saturday,'%d %B %Y')
    date_text.text = get_next_weekday(datetime.today().strftime('%d %B %Y'), 5)
    church_name.text = "Reformed Evangelical Church Singapore"

def add_first_song_title_slide():
    print("First Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_1['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = song_1['author']
    number_text.text = "1"

def add_second_song_title_slide():
    print("Second Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_2['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = song_2['author']
    number_text.text = "2"

def add_third_song_title_slide():
    print("Third Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_3['title']
    number_text = slide.placeholders[10]
    author_text = slide.placeholders[11]
    author_text.text = song_3['author']
    number_text.text = "3"
    
def add_first_song_lyrics_slide(verse_number):
    print("First Lyric Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_1['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = song_1['author']
    print(verse_number)
    lyrics_text.text = song_1['lyrics'][verse_number][0]
    number_text.text = generate_verse_number(song_1['verse_number'][verse_number])
    for i in range(1, len(song_1['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_1['lyrics'][verse_number][i]

def add_second_song_lyrics_slide(verse_number):
    print("First Lyric Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_2['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = song_2['author']
    number_text.text = generate_verse_number(song_2['verse_number'][verse_number])
    lyrics_text.text = song_2['lyrics'][verse_number][0]
    for i in range(1, len(song_2['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_2['lyrics'][verse_number][i]

def add_third_song_lyrics_slide(verse_number):
    print("First Lyric Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.title
    for shape in slide.placeholders:
     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    title.text = song_3['title']
    author_text = slide.placeholders[11]
    number_text = slide.placeholders[13]
    lyrics_text = slide.placeholders[12].text_frame
    author_text.text = song_3['author']
    number_text.text = generate_verse_number(song_3['verse_number'][verse_number])
    lyrics_text.text = song_3['lyrics'][verse_number][0]
    for i in range(1, len(song_3['lyrics'][verse_number])):
        p = lyrics_text.add_paragraph()
        p.text = song_3['lyrics'][verse_number][i]    

def add_black_slide():
    prs.slides.add_slide(prs.slide_layouts[7])

def add_announcements():
    prs.slides.add_slide(prs.slide_layouts[8])
    prs.slides.add_slide(prs.slide_layouts[9])
    birthday = prs.slides.add_slide(prs.slide_layouts[10])
    birthday.shapes.title.text = "Birthday Song"
    prs.slides.add_slide(prs.slide_layouts[11])

def generate_first_song():
    add_first_song_title_slide()
    for i in range(0, len(song_1['verse_number'])):
        add_first_song_lyrics_slide(i)
    add_black_slide()

def generate_second_song():
    add_second_song_title_slide()
    for i in range(0, len(song_2['verse_number'])):
        add_second_song_lyrics_slide(i)
    add_black_slide()

def generate_third_song():
    add_third_song_title_slide()
    for i in range(0, len(song_3['verse_number'])):
        add_third_song_lyrics_slide(i)
    add_black_slide()
id_to_song()
add_welcome_slide()
generate_first_song()
generate_second_song()
generate_third_song() 
add_announcements() 

prs.save(output_presentation)
